import streamlit as st
from pptx import Presentation
import torch
from transformers import AutoTokenizer, AutoModel
import io
import json
import pandas as pd
import re
import requests

# -------------------------
# Device Setup (CPU only)
# -------------------------
device = torch.device("cpu")

# -------------------------
# Embedding Model Setup
# -------------------------
def load_model():
    tokenizer = AutoTokenizer.from_pretrained("sentence-transformers/all-MiniLM-L6-v2")
    model = AutoModel.from_pretrained("sentence-transformers/all-MiniLM-L6-v2").to(device)
    return tokenizer, model

def get_embeddings(texts, tokenizer, model):
    inputs = tokenizer(texts, padding=True, truncation=True, return_tensors="pt").to(device)
    with torch.no_grad():
        embeddings = model(**inputs).last_hidden_state.mean(dim=1)
    return embeddings  # stays on CPU

# -------------------------
# Extract text from PPTX
# -------------------------
def extract_text_from_pptx(file):
    prs = Presentation(file)
    all_text = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text.strip())
        if slide_text:
            all_text.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
    return "\n\n".join(all_text)

# -------------------------
# Groq API call
# -------------------------
def call_groq_api(prompt, groq_api_key, model="llama-3.1-8b-instant", temperature=0.4, seed=None):
    url = "https://api.groq.com/openai/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {groq_api_key}",
        "Content-Type": "application/json"
    }
    data = {
        "model": model,
        "messages": [{"role": "user", "content": prompt}],
        "temperature": temperature
    }
    if seed is not None:
        data["seed"] = seed

    response = requests.post(url, headers=headers, data=json.dumps(data))
    if response.status_code == 200:
        return response.json()["choices"][0]["message"]["content"]
    else:
        raise Exception(f"Groq API Error {response.status_code}: {response.text}")

# -------------------------
# Generate MCQs in Batches (Groq)
# -------------------------
def generate_mcqs_in_batches(text, total_questions, batch_size, seed, groq_api_key, model_choice):
    all_mcqs = []
    num_batches = (total_questions + batch_size - 1) // batch_size  # ceiling division

    text_chunks = [text[i:i+4000] for i in range(0, len(text), 4000)]
    if not text_chunks:
        text_chunks = [text]

    for i in range(num_batches):
        questions_to_generate = min(batch_size, total_questions - len(all_mcqs))
        st.info(f"🔄 Generating batch {i+1}/{num_batches} ({questions_to_generate} MCQs)...")

        content = text_chunks[i % len(text_chunks)]

        prompt = f"""
        Based on the following content, generate {questions_to_generate} exam-quality multiple-choice questions (MCQs).
        Each MCQ must include:
        - A clear question
        - Four options (A, B, C, D)
        - The correct answer (just the letter)
        - Number the MCQs as Q1, Q2, etc.

        Return ONLY a JSON array like this:
        [
          {{
            "question": "Q1. Question text",
            "options": ["Option A", "Option B", "Option C", "Option D"],
            "answer": "A"
          }}
        ]

        Content:
        {content}
        """

        try:
            output = call_groq_api(prompt, groq_api_key, model=model_choice, seed=seed)
            json_match = re.search(r'\[.*\]', output, re.DOTALL)
            if json_match:
                mcqs = json.loads(json_match.group(0))
                all_mcqs.extend(mcqs)
            else:
                st.warning(f"⚠️ Batch {i+1} did not return valid JSON.")
        except Exception as e:
            st.error(f"❌ Error in batch {i+1}: {e}")

    # Ensure exact required number
    if len(all_mcqs) > total_questions:
        all_mcqs = all_mcqs[:total_questions]
    elif len(all_mcqs) < total_questions:
        needed = total_questions - len(all_mcqs)
        filler = all_mcqs[:needed] if all_mcqs else []
        all_mcqs.extend(filler)

    # Re-number to maintain consistency
    for i, mcq in enumerate(all_mcqs, start=1):
        qtext = mcq.get("question", "")
        if not qtext.strip().lower().startswith("q"):
            mcq["question"] = f"Q{i}. {qtext}"
        else:
            mcq["question"] = f"Q{i}. {qtext.split('.', 1)[-1].strip()}"

    return all_mcqs

# -------------------------
# Convert MCQs to Excel
# -------------------------
def mcqs_to_excel(mcqs):
    data = []
    for i, mcq in enumerate(mcqs, start=1):
        question = mcq.get("question", "")
        options = mcq.get("options", [])
        answer = mcq.get("answer", "")

        while len(options) < 4:
            options.append("")
        options = options[:4]

        data.append([i, question, options[0], options[1], options[2], options[3], answer])

    df = pd.DataFrame(data, columns=["Q.No", "Question", "Option A", "Option B", "Option C", "Option D", "Answer"])

    output = io.BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name="MCQs")
    output.seek(0)
    return output

# -------------------------
# Convert MCQs to JSON
# -------------------------
def mcqs_to_json(mcqs):
    return io.BytesIO(json.dumps(mcqs, indent=2).encode("utf-8"))

# -------------------------
# Streamlit UI
# -------------------------
st.title("📊 PPTX to MCQ Generator (Groq API)")

uploaded_file = st.file_uploader("Upload a PowerPoint (.pptx) file", type=["pptx"])

# Sidebar for API Key and Model Choice
st.sidebar.header("🔑 Groq Settings")
groq_api_key = st.sidebar.text_input("Enter your Groq API Key", type="password")
model_choice = st.sidebar.selectbox(
    "Choose Groq Model",
    ["llama-3.1-8b-instant", "llama-3.3-70b-versatile"],
    index=0
)

if uploaded_file is not None:
    if not groq_api_key:
        st.warning("⚠️ Please enter your Groq API Key in the sidebar to continue.")
    else:
        with st.spinner("Extracting text..."):
            ppt_text = extract_text_from_pptx(uploaded_file)
            st.subheader("📄 Extracted Text:")
            st.text_area("Text Output", ppt_text, height=300)

            text_file = io.BytesIO(ppt_text.encode("utf-8"))
            st.download_button("⬇ Download Extracted Text", text_file, "extracted_ppt_text.txt", "text/plain")

            # MCQ Generation options
            st.subheader("MCQ Generation Options")
            num_mcqs = st.slider("Select number of MCQs to generate", 10, 100, step=10, value=20)
            seed_value = st.number_input("Seed value (optional)", min_value=0, step=1, value=0)

            if st.button("Generate MCQs"):
                with st.spinner(f"Generating {num_mcqs} MCQs in batches..."):
                    mcqs = generate_mcqs_in_batches(
                        ppt_text,
                        total_questions=num_mcqs,
                        batch_size=10,
                        seed=seed_value if seed_value > 0 else None,
                        groq_api_key=groq_api_key,
                        model_choice=model_choice
                    )

                if not mcqs:
                    st.error("⚠️ MCQ generation failed.")
                else:
                    st.success(f"✅ Successfully generated {len(mcqs)} MCQs!")
                    st.subheader("📝 Generated MCQs Preview:")
                    st.json(mcqs[:5])  # Preview first 5

                    excel_file = mcqs_to_excel(mcqs)
                    json_file = mcqs_to_json(mcqs)

                    st.download_button(
                        label=f"⬇ Download {len(mcqs)} MCQs (Excel)",
                        data=excel_file,
                        file_name=f"generated_{len(mcqs)}_mcqs.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                    )
                    st.download_button(
                        label=f"⬇ Download {len(mcqs)} MCQs (JSON)",
                        data=json_file,
                        file_name=f"generated_{len(mcqs)}_mcqs.json",
                        mime="application/json"
                    )
